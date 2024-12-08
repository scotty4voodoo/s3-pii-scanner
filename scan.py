import json
import re
import os
from datetime import datetime
from pptx import Presentation
import docx
import PyPDF2
import pandas as pd

def write_scan_results_to_file():
    # Create scanning-result directory if it doesn't exist
    if not os.path.exists('scanning-result'):
        os.makedirs('scanning-result')

    # Generate filename with current date
    output_file = os.path.join('scanning-result', f"{datetime.now().strftime('%Y-%m-%d-%s')}.txt")

    # Redirect stdout to file
    with open(output_file, 'w') as f:
        # Load regex patterns from JSON file
        with open('patterns.json', 'r') as patterns_file:
            patterns = json.load(patterns_file)

# Walk through downloaded files
        for root, dirs, files in os.walk('temp_scan_files'):
            for file in files:
                file_path = os.path.join(root, file)
                
                # Process binary files
                if file_path.endswith((".ppt", ".pptx", ".doc", ".docx", ".xls", ".xlsx", ".pdf")):
                    try:
                        # Import appropriate library based on file extension
                        if file_path.endswith((".ppt", ".pptx")):
                            prs = Presentation(file_path)
                            content = ""
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        content += shape.text
                                        
                        elif file_path.endswith((".doc", ".docx")):
                            doc = docx.Document(file_path)
                            content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                            
                        elif file_path.endswith((".xls", ".xlsx")):
                            
                            df = pd.read_excel(file_path)
                            content = df.to_string()
                            
                        elif file_path.endswith(".pdf"):
                            reader = PyPDF2.PdfReader(file_path)
                            content = ""
                            for page in reader.pages:
                                content += page.extract_text()
                        # Check patterns
                        for pattern_name, pattern in patterns.items():
                            match = re.search(pattern, content)
                            if match:
                                f.write(f"{file_path} : {match.group()} 에서 {pattern_name} 을 탐지했습니다.\n")
                    except Exception as e:
                        f.write(f"{file_path}: {str(e)} 파일을 스캐닝 하는 과정에서 에러가 발생했습니다.\n")
                else:
                    try:
                        with open(file_path, 'r') as scan_file:
                            content = scan_file.read()

                        # Check patterns
                        # Break out of pattern checking once a match is foun
                        for pattern_name, pattern in patterns.items():
                            match = re.search(pattern, content)
                            if match:
                                f.write(f"{file_path} : {match.group()} 에서 {pattern_name} 을 탐지했습니다.\n")

                    except Exception as e:
                        f.write(f"{file_path}: {str(e)} 파일을 스캐닝 하는 과정에서 에러가 발생했습니다.\n")

        print(f"스캔 결과가 {output_file} 에 기록되었습니다.")
        return output_file
